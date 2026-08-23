"""
PPT operacional: tratamento FPD / Qualidade no ClickUp.
Identidade visual alinhada ao Design System ClickUp.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
DANGER = RGBColor(0xDC, 0x26, 0x26)
WARNING = RGBColor(0xD9, 0x77, 0x06)
BG = RGBColor(0xF7, 0xF9, 0xFC)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x23, 0x32)
TEXT_MUTED = RGBColor(0x5B, 0x6B, 0x7D)
FOOTER = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TIP_BG = RGBColor(0xE0, 0xF2, 0xFE)
WARN_BG = RGBColor(0xFE, 0xF3, 0xC7)
OK_BG = RGBColor(0xD1, 0xFA, 0xE5)
TEAL = RGBColor(0x0F, 0x76, 0x6E)

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "static" / "logo.png"
OUT = ROOT / "docs" / "Qualidade_Tratamento_FPD_Passo_a_Passo.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    return shape


def _add_round_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold=bold, color=color)
    return box


def _bullets(slide, left, top, width, height, lines: list[str], size: int = 15) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("• " if not line.startswith("•") else "") + line
        _set_run_font(run, size, color=TEXT)


def _footer(slide, page: int, total: int) -> None:
    _add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), FOOTER)
    _textbox(
        slide,
        Inches(0.4),
        Inches(7.18),
        Inches(8),
        Inches(0.28),
        "ClickUp · Qualidade / Tratamento FPD · Uso operacional",
        size=11,
        color=WHITE,
    )
    _textbox(
        slide,
        Inches(11.2),
        Inches(7.18),
        Inches(1.8),
        Inches(0.28),
        f"{page} / {total}",
        size=11,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )


def _header(slide, title: str, subtitle: str = "") -> None:
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), FOOTER)
    if LOGO.exists():
        try:
            slide.shapes.add_picture(str(LOGO), Inches(0.35), Inches(0.22), height=Inches(0.6))
        except Exception:
            pass
    _textbox(slide, Inches(1.5), Inches(0.18), Inches(10.5), Inches(0.45), title, size=24, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, Inches(1.5), Inches(0.58), Inches(10.5), Inches(0.35), subtitle, size=13, color=ACCENT)


def _base(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG)
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slides_meta: list[tuple] = []

    # 1 capa
    s = _base(prs)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, FOOTER)
    _textbox(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1), "Tratamento FPD — Qualidade", size=36, bold=True, color=WHITE)
    _textbox(
        s,
        Inches(0.8),
        Inches(3.3),
        Inches(11.5),
        Inches(1),
        "Passo a passo operacional no ClickUp\nO que é manual · o que é automático · WhatsApp oficial",
        size=18,
        color=ACCENT,
    )
    slides_meta.append(s)

    # 2 objetivo
    s = _base(prs)
    _header(s, "O que é o Tratamento FPD", "Módulo Qualidade · fila de atrasados / em aberto / pagas")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(11.5),
        Inches(5),
        [
            "Acompanhar a 1ª fatura (FPD) e faturas seguintes (SPD/TPD) após a instalação.",
            "Cobrar clientes em atraso com WhatsApp oficial (templates Meta) ou e-mail.",
            "Registrar status de tratamento, promessa de pagamento, PIX/barras e ligações.",
            "Conferir o que a planilha da operadora confirma versus o que o time informou.",
            "Histórico de contatos: envios, ligações e respostas do cliente (botões/texto).",
        ],
        size=17,
    )
    slides_meta.append(s)

    # 3 jornada
    s = _base(prs)
    _header(s, "Jornada do processo (visão geral)", "Da planilha FPD até a cobrança e o retorno do cliente")
    steps = [
        ("1", "Importar FPD", "Planilha da operadora\npreenche status/venc."),
        ("2", "Fila Tratamento", "Atrasados / aberto\n/ pagas"),
        ("3", "Completar dados", "Valor, PIX, barras\nvia Nio ou manual"),
        ("4", "Contatar", "WhatsApp · e-mail\n· ligação Sonax"),
        ("5", "Registrar", "Status trat. + hist.\n+ resposta cliente"),
    ]
    x0 = 0.45
    for i, (n, t, d) in enumerate(steps):
        left = Inches(x0 + i * 2.55)
        _add_round_rect(s, left, Inches(2.0), Inches(2.35), Inches(3.4), SURFACE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(2.25), Inches(0.7), Inches(0.7))
        _fill_solid(circ, ACCENT)
        _textbox(s, left + Inches(0.8), Inches(2.35), Inches(0.7), Inches(0.5), n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(s, left + Inches(0.15), Inches(3.15), Inches(2.05), Inches(0.5), t, size=15, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        _textbox(s, left + Inches(0.15), Inches(3.7), Inches(2.05), Inches(1.4), d, size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    slides_meta.append(s)

    # 4 passo importar
    s = _base(prs)
    _header(s, "Passo 1 — Importar a planilha FPD", "Central de Importações / fluxo FPD·SPD·TPD")
    _add_round_rect(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(5.2), SURFACE)
    _textbox(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.4), "O que a importação faz", size=16, bold=True, color=TEAL)
    _bullets(
        s,
        Inches(0.75),
        Inches(2.1),
        Inches(5.5),
        Inches(4.2),
        [
            "Lê NR_ORDEM e casa com a venda/contrato no CRM.",
            "Atualiza faturas 1/2/3 (FPD/SPD/TPD).",
            "Preenche/atualiza status, vencimento e campos espelho no contrato.",
            "Marca O.S. sem CRM em “Faltam no CRM”.",
            "Se o BO já marcou pago no tratamento: confirma ou marca divergente.",
        ],
        size=14,
    )
    _add_round_rect(s, Inches(6.8), Inches(1.35), Inches(5.9), Inches(5.2), TIP_BG)
    _textbox(s, Inches(7.05), Inches(1.55), Inches(5.4), Inches(0.4), "Dica operacional", size=16, bold=True, color=ACCENT)
    _bullets(
        s,
        Inches(7.05),
        Inches(2.1),
        Inches(5.4),
        Inches(4.2),
        [
            "Sem importação recente a fila fica desalinhada da operadora.",
            "Valor da fatura nem sempre vem completo na planilha — confira na Nio.",
            "Órfãos / sem CPF bloqueiam cobrança WhatsApp.",
        ],
        size=14,
    )
    slides_meta.append(s)

    # 5 passo fila
    s = _base(prs)
    _header(s, "Passo 2 — Abrir a fila de Tratamento", "Qualidade → aba Tratamento")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(11.5),
        Inches(5.2),
        [
            "Escolha o mês (vencimento da 1ª fatura) e a lente “Venc. 1ª / Instalação”.",
            "Filas: ATRASADOS · EM ABERTO · PAGAS · TOTAL.",
            "Filtros: atraso (dias), faturas pagas, status de tratamento, busca O.S./cliente/CPF.",
            "CPF aparece abaixo do nome; coluna Valor 1ª (time) ou Bônus (Diretoria/Admin).",
            "Badges Contato: WhatsApp · e-mail · ligação · resposta — clique abre o histórico.",
        ],
        size=16,
    )
    slides_meta.append(s)

    # 6 completar dados
    s = _base(prs)
    _header(s, "Passo 3 — Completar valor, vencimento e 2ª via", "Sem valor válido o envio WhatsApp fica bloqueado")
    _add_round_rect(s, Inches(0.5), Inches(1.35), Inches(4.0), Inches(5.2), SURFACE)
    _textbox(s, Inches(0.7), Inches(1.55), Inches(3.6), Inches(0.4), "Automático", size=16, bold=True, color=SUCCESS)
    _bullets(
        s,
        Inches(0.7),
        Inches(2.1),
        Inches(3.6),
        Inches(4.2),
        [
            "Importação FPD: status + vencimento da planilha.",
            "Job 00:05: busca faturas Nio (histórico).",
            "Buscar Nio na tela: valor, PIX, barras, PDF.",
        ],
        size=13,
    )
    _add_round_rect(s, Inches(4.7), Inches(1.35), Inches(4.0), Inches(5.2), SURFACE)
    _textbox(s, Inches(4.9), Inches(1.55), Inches(3.6), Inches(0.4), "Manual na Qualidade", size=16, bold=True, color=WARNING)
    _bullets(
        s,
        Inches(4.9),
        Inches(2.1),
        Inches(3.6),
        Inches(4.2),
        [
            "Abrir faturas (ícone recibo).",
            "Editar valor, vencimento, promessa.",
            "Colar PIX / código de barras.",
            "Salvar faturas.",
            "Status tratado aguarda próximo FPD.",
        ],
        size=13,
    )
    _add_round_rect(s, Inches(8.9), Inches(1.35), Inches(3.9), Inches(5.2), WARN_BG)
    _textbox(s, Inches(9.1), Inches(1.55), Inches(3.5), Inches(0.4), "Bloqueio de envio", size=16, bold=True, color=DANGER)
    _bullets(
        s,
        Inches(9.1),
        Inches(2.1),
        Inches(3.5),
        Inches(4.2),
        [
            "Valor R$ 0,00 ou vazio → não envia.",
            "Sem vencimento → não envia.",
            "Sem CPF / órfão → não envia.",
            "Corrija a fatura e tente de novo.",
        ],
        size=13,
    )
    slides_meta.append(s)

    # 7 contato
    s = _base(prs)
    _header(s, "Passo 4 — Contatar o cliente", "WhatsApp oficial (Número B) · e-mail · ligação Sonax")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(11.5),
        Inches(5.2),
        [
            "WhatsApp: botão verde → confira telefone, valor e vencimento → Enviar (veja o status do botão).",
            "Template Meta automático conforme dias: lembrete (antes), vencida, recorrente.",
            "E-mail: roteiro de 2ª via quando houver e-mail no cadastro.",
            "Ligar: botão telefone → escolha ramal Sonax → a ligação fica no histórico.",
            "Tudo fica nos badges Contato; respostas do cliente (botão/texto) no modal de histórico.",
        ],
        size=16,
    )
    slides_meta.append(s)

    # 8 respostas
    s = _base(prs)
    _header(s, "Passo 5 — O que o cliente responde no WhatsApp", "Botões do template Meta + texto livre")
    _add_round_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.6), TIP_BG)
    _textbox(
        s,
        Inches(0.75),
        Inches(1.6),
        Inches(11.8),
        Inches(1.2),
        "Como o sistema reconhece o botão: o WhatsApp envia o texto do Quick Reply "
        "(ex.: QUERO A 2A VIA). O webhook normaliza e classifica esse texto.",
        size=15,
        color=TEXT,
    )
    _bullets(
        s,
        Inches(0.7),
        Inches(3.3),
        Inches(11.5),
        Inches(3.3),
        [
            "QUERO A 2A VIA → envia roteiro com PIX/barras e grava resposta.",
            "JA PAGUEI → pede comprovante e grava resposta.",
            "FALAR COM SUPORTE → avisa que um especialista falará e grava resposta.",
            "Texto livre (depois de uma cobrança nos últimos 14 dias) também é gravado.",
            "Ver tudo: clique nos badges Contato na linha do contrato.",
        ],
        size=15,
    )
    slides_meta.append(s)

    # 9 automação cobrança
    s = _base(prs)
    _header(s, "Cobrança automática — quando dispara?", "Scheduler do site-clickup (America/Sao_Paulo)")
    _add_round_rect(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.2), SURFACE)
    _bullets(
        s,
        Inches(0.8),
        Inches(1.6),
        Inches(11.7),
        Inches(4.7),
        [
            "Todo dia às 10:00 — job enviar_templates_cobranca_nio (se WHATSAPP_USE_NIO_TEMPLATES=true).",
            "D−5 (5 dias antes do vencimento) → template lembrete.",
            "D+5 (5 dias após o vencimento) → template fatura vencida.",
            "D+12, D+19, D+26… (a cada 7 dias após D+5) → template recorrente.",
            "Só faturas NAO_PAGO / ATRASADO / AGUARDANDO, contrato tratável, e sem WhatsApp sucesso no mesmo dia.",
            "Limite padrão por execução: 80 envios. Envio manual na tela NÃO depende desse horário.",
        ],
        size=15,
    )
    slides_meta.append(s)

    # 10 preenchimento automático
    s = _base(prs)
    _header(s, "Preenchimento automático de dados — o que existe hoje", "Valor · vencimento · status · PIX")
    rows = [
        ("Importação FPD", "Status, vencimento, espelho no contrato; match O.S. × CRM", "Ao importar a planilha"),
        ("Busca Nio 00:05", "Varredura diária de faturas (comando automático)", "Todo dia 00:05"),
        ("Botão Buscar Nio", "Valor, PIX, barras, PDF da opção escolhida", "Manual na Qualidade"),
        ("Salvar faturas", "Edits manuais + status aguardando FPD", "Manual na Qualidade"),
        ("Cobrança 10:00", "Só dispara template se valor/venc. já válidos", "Automático diário"),
    ]
    y = 1.35
    _add_round_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.55), TEAL)
    for i, h in enumerate(["Fonte", "O que preenche", "Quando"]):
        _textbox(s, Inches(0.7 + i * 4.0), y + Inches(0.1), Inches(3.8), Inches(0.4), h, size=14, bold=True, color=WHITE)
    for idx, (a, b, c) in enumerate(rows):
        yy = Inches(2.05 + idx * 0.85)
        bg = SURFACE if idx % 2 == 0 else TIP_BG
        _add_round_rect(s, Inches(0.5), yy, Inches(12.3), Inches(0.75), bg)
        _textbox(s, Inches(0.7), yy + Inches(0.15), Inches(3.8), Inches(0.5), a, size=13, bold=True, color=TEXT)
        _textbox(s, Inches(4.7), yy + Inches(0.15), Inches(3.8), Inches(0.5), b, size=12, color=TEXT)
        _textbox(s, Inches(8.7), yy + Inches(0.15), Inches(3.8), Inches(0.5), c, size=12, color=TEXT_MUTED)
    slides_meta.append(s)

    # 11 checklist diário
    s = _base(prs)
    _header(s, "Checklist diário do analista", "Rotina sugerida no Tratamento")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.4),
        Inches(11.5),
        Inches(5.2),
        [
            "1. Confirmar se a última planilha FPD foi importada.",
            "2. Abrir fila ATRASADOS do mês e filtrar por faixa de atraso.",
            "3. Em cada caso: abrir faturas → validar valor/venc./PIX (Buscar Nio se faltar).",
            "4. Enviar WhatsApp / e-mail ou ligar; acompanhar badges Contato.",
            "5. Atualizar Status tratamento e promessa de pagamento.",
            "6. Revisar respostas do cliente no histórico (botão ou texto).",
            "7. No próximo FPD, conferir badges Aguard. FPD / OK FPD / Divergente.",
        ],
        size=16,
    )
    slides_meta.append(s)

    # 12 encerramento
    s = _base(prs)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, FOOTER)
    _textbox(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.7), "Pronto para operar", size=32, bold=True, color=WHITE)
    _textbox(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11.5),
        Inches(1.5),
        "Automático: importação FPD · busca Nio 00:05 · templates de cobrança 10:00\n"
        "Manual: completar valor/PIX, contatar, registrar status e ler respostas\n"
        "Dúvidas técnicas: híbrido Z-API (equipe) + WhatsAtende B (cliente)",
        size=16,
        color=ACCENT,
    )
    slides_meta.append(s)

    total = len(slides_meta)
    for i, slide in enumerate(slides_meta, start=1):
        if i == 1 or i == total:
            continue
        _footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"PPT gerado: {path}")
