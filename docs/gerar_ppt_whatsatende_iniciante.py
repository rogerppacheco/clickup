"""
Gera PPT iniciante: como começar a usar a WhatsAtende e integrar no site-clickup.
Identidade visual alinhada ao Design System ClickUp v17.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Cores ClickUp v17 ---
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT_HOVER = RGBColor(0x02, 0x84, 0xC7)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
DANGER = RGBColor(0xDC, 0x26, 0x26)
WARNING = RGBColor(0xD9, 0x77, 0x06)
BG = RGBColor(0xF7, 0xF9, 0xFC)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x23, 0x32)
TEXT_MUTED = RGBColor(0x5B, 0x6B, 0x7D)
FOOTER = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TIP_BG = RGBColor(0xE0, 0xF2, 0xFE)
WARN_BG = RGBColor(0xFE, 0xF3, 0xC7)
OK_BG = RGBColor(0xD1, 0xFA, 0xE5)
TEAL = RGBColor(0x0F, 0x76, 0x6E)

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "static" / "logo.png"
OUT = ROOT / "docs" / "WhatsAtende_Guia_Iniciante_Record_PAP.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT, name: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    return shape


def _add_round_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set("anchor", "t" if anchor == MSO_ANCHOR.TOP else "ctr")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold, color)
    return box


def _add_paragraph(
    tf,
    text: str,
    size: int = 13,
    bold: bool = False,
    color: RGBColor = TEXT,
    space_before: int = 4,
    space_after: int = 2,
    align=PP_ALIGN.LEFT,
):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold, color)
    return p


def _base_slide(prs: Presentation, with_footer: bool = True):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), TEAL)
    if with_footer:
        _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), FOOTER)
        _textbox(
            slide,
            Inches(0.4),
            Inches(7.18),
            Inches(9),
            Inches(0.28),
            "ClickUp  ·  WhatsAtende — guia iniciante  ·  Uso interno",
            size=10,
            color=RGBColor(0x94, 0xA3, 0xB8),
        )
        _textbox(
            slide,
            Inches(10.5),
            Inches(7.18),
            Inches(2.5),
            Inches(0.28),
            "Confidencial",
            size=10,
            color=RGBColor(0x94, 0xA3, 0xB8),
            align=PP_ALIGN.RIGHT,
        )
    return slide


def _header(slide, title: str, subtitle: str = "", logo: bool = True):
    if logo and LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.4), Inches(0.22), height=Inches(0.55))
    left = Inches(2.2) if logo and LOGO.exists() else Inches(0.4)
    _textbox(slide, left, Inches(0.22), Inches(10), Inches(0.4), title, size=24, bold=True, color=TEXT)
    if subtitle:
        _textbox(slide, left, Inches(0.62), Inches(10.5), Inches(0.3), subtitle, size=12, color=TEXT_MUTED)


def _chip(slide, left, top, width, height, text: str, bg: RGBColor, fg: RGBColor):
    _add_round_rect(slide, left, top, width, height, bg)
    _textbox(
        slide,
        left,
        top + Inches(0.05),
        width,
        height - Inches(0.05),
        text,
        size=11,
        bold=True,
        color=fg,
        align=PP_ALIGN.CENTER,
    )


def _card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    lines: list[str],
    title_color: RGBColor = ACCENT_HOVER,
    accent_bar: RGBColor | None = TEAL,
):
    _add_round_rect(slide, left, top, width, height, SURFACE)
    if accent_bar:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        _fill_solid(bar, accent_bar)
    box = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.12),
        width - Inches(0.35),
        height - Inches(0.2),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run_font(run, 13, True, title_color)
    for line in lines:
        _add_paragraph(tf, line, size=12, color=TEXT, space_before=3, space_after=1)


def _screen_mock(slide, left, top, width, height, title: str, lines: list[str], tip: str = ""):
    """Simula uma 'tela' do painel WhatsAtende."""
    _add_round_rect(slide, left, top, width, height, FOOTER)
    _add_rect(slide, left, top, width, Inches(0.4), TEAL)
    _textbox(
        slide,
        left + Inches(0.15),
        top + Inches(0.08),
        width - Inches(0.3),
        Inches(0.28),
        title,
        size=12,
        bold=True,
        color=WHITE,
    )
    box = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.55),
        width - Inches(0.4),
        height - Inches(0.7),
    )
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        _set_run_font(run, 11, False, RGBColor(0xE2, 0xE8, 0xF0))
    if tip:
        _textbox(
            slide,
            left + Inches(0.15),
            top + height - Inches(0.38),
            width - Inches(0.3),
            Inches(0.3),
            tip,
            size=10,
            bold=True,
            color=RGBColor(0x67, 0xE8, 0xF9),
        )


def _step_row(slide, left, top, width, num: str, title: str, detail: str):
    _add_round_rect(slide, left, top, Inches(0.45), Inches(0.45), TEAL)
    _textbox(
        slide,
        left,
        top + Inches(0.08),
        Inches(0.45),
        Inches(0.35),
        num,
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _textbox(slide, left + Inches(0.6), top, width - Inches(0.6), Inches(0.28), title, size=14, bold=True, color=TEXT)
    _textbox(
        slide,
        left + Inches(0.6),
        top + Inches(0.28),
        width - Inches(0.6),
        Inches(0.35),
        detail,
        size=12,
        color=TEXT_MUTED,
    )


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def slide_capa(prs: Presentation) -> None:
    slide = _base_slide(prs, with_footer=False)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, FOOTER)
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, TEAL)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.7), Inches(1.4), height=Inches(0.7))
    _textbox(
        slide,
        Inches(0.7),
        Inches(2.4),
        Inches(11),
        Inches(0.7),
        "WhatsAtende — Guia para iniciantes",
        size=36,
        bold=True,
        color=WHITE,
    )
    _textbox(
        slide,
        Inches(0.7),
        Inches(3.2),
        Inches(11),
        Inches(0.5),
        "Do primeiro login até funcionar com o site-clickup (bot ClickUp)",
        size=18,
        color=RGBColor(0x94, 0xA3, 0xB8),
    )
    _chip(slide, Inches(0.7), Inches(4.2), Inches(2.2), Inches(0.4), "Passo a passo", TEAL, WHITE)
    _chip(slide, Inches(3.1), Inches(4.2), Inches(2.4), Inches(0.4), "Telas do painel", ACCENT, WHITE)
    _chip(slide, Inches(5.7), Inches(4.2), Inches(2.8), Inches(0.4), "Integração projeto", SUCCESS, WHITE)
    _textbox(
        slide,
        Inches(0.7),
        Inches(6.5),
        Inches(11),
        Inches(0.4),
        "ClickUp LTDA  ·  Instância app14.whatsatende.com.br  ·  Uso interno",
        size=12,
        color=RGBColor(0x94, 0xA3, 0xB8),
    )


def slide_objetivo(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Para que vamos usar a WhatsAtende?", "Dois números, duas finalidades")
    _card(
        slide,
        Inches(0.4),
        Inches(1.2),
        Inches(6.0),
        Inches(5.4),
        "Número A — Bot da equipe comercial",
        [
            "Automações internas (comandos do bot).",
            "Comportamento parecido com Z-API (mensagem livre).",
            "Conexão recomendada: QR Code (WhatsApp Web).",
            "Atenção: não é API oficial da Meta.",
            "Risco de desconexão/bloqueio deve ficar no contrato.",
            "",
            "No site-clickup: provedor WhatsAtende (quando ativo).",
        ],
        title_color=TEAL,
        accent_bar=TEAL,
    )
    _card(
        slide,
        Inches(6.7),
        Inches(1.2),
        Inches(6.0),
        Inches(5.4),
        "Número B — Cliente final",
        [
            "Pós-venda, instalação, fatura, cobrança.",
            "Conexão recomendada: Cloud API oficial (Meta).",
            "Fora da janela de 24h: só template aprovado.",
            "Dentro das 24h: texto livre permitido.",
            "WABA no Business Manager da ClickUp.",
            "",
            "Melhor governança e menor risco comercial.",
        ],
        title_color=ACCENT_HOVER,
        accent_bar=ACCENT,
    )


def slide_mapa(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Mapa do caminho (checklist)", "Siga nesta ordem — não pule etapas")
    steps = [
        ("1", "Acessar o painel", "Login em app14.whatsatende.com.br"),
        ("2", "Criar a conexão", "Menu Conexões → Novo (QR ou Oficial)"),
        ("3", "Gerar o token da API", "Editar conexão → campo Token → copiar"),
        ("4", "Anotar o ID da conexão", "whatsappId (necessário no .env / Railway)"),
        ("5", "Configurar webhook", "URL do site-clickup (quando o suporte confirmar)"),
        ("6", "Variáveis no servidor", "WHATSATENDE_TOKEN + WHATSATENDE_WHATSAPP_ID"),
        ("7", "Ativar no painel ClickUp", "Configuração WhatsApp → provedor WhatsAtende"),
        ("8", "Testar envio + recebimento", "Texto, PDF e mensagem de teste no bot"),
    ]
    y = Inches(1.15)
    for i, (num, title, detail) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = Inches(0.5) if col == 0 else Inches(6.9)
        top = y + Inches(row * 1.35)
        _add_round_rect(slide, left, top, Inches(5.9), Inches(1.15), SURFACE)
        _step_row(slide, left + Inches(0.2), top + Inches(0.25), Inches(5.4), num, title, detail)


def slide_login(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 1 — Entrar no painel", "Tela de login")
    _screen_mock(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(6.2),
        Inches(5.4),
        "WhatsAtende · Login",
        [
            "URL: https://app14.whatsatende.com.br/login",
            "",
            "Campo Email  →  recordpap@gmail.com",
            "Campo Senha  →  (sua senha atual)",
            "",
            "[ Entrar ]",
            "",
            "Após o login você cai em Atendimentos.",
            "Menu lateral: Gerência / Administração / API.",
        ],
        tip="Dica: troque a senha inicial assim que possível.",
    )
    _card(
        slide,
        Inches(7.0),
        Inches(1.2),
        Inches(5.7),
        Inches(5.4),
        "O que você vai ver depois do login",
        [
            "• Atendimentos (tickets)",
            "• Contatos",
            "• Campanhas / Automações / Fluxos",
            "• API → Documentação / Teste / Logs",
            "• Conexões ← aqui começa a integração",
            "• Templates WABA (API oficial)",
            "• Configurações",
            "",
            "Conta: CLICKUP LTDA",
            "Plano com External API liberada.",
            "Até 2 conexões no plano atual.",
        ],
        title_color=TEAL,
    )


def slide_conexao_qr(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 2 — Criar conexão (Número A / QR)", "Menu Administração → Conexões")
    _screen_mock(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(6.3),
        Inches(5.5),
        "Conexões · Nova conexão",
        [
            "1. Clique em Novo",
            "2. Nome: ex. BOT_EQUIPE_COMERCIAL",
            "3. Canal: WhatsApp Web (Baileys / QR)",
            "4. Salve a conexão",
            "5. Clique em conectar / gerar QR",
            "6. Abra WhatsApp no celular →",
            "   Aparelhos conectados → Vincular",
            "7. Escaneie o QR até status Conectado",
            "",
            "Anote o ID da conexão (whatsappId)",
            "que aparece na edição/card.",
        ],
        tip="Use um chip/número dedicado ao bot interno.",
    )
    _card(
        slide,
        Inches(7.0),
        Inches(1.15),
        Inches(5.7),
        Inches(2.5),
        "Checklist desta tela",
        [
            "☐ Conexão criada com nome claro",
            "☐ Status = Conectado (não “Aguardando”)",
            "☐ ID da conexão copiado",
            "☐ Celular com internet estável",
        ],
        title_color=SUCCESS,
        accent_bar=SUCCESS,
    )
    _card(
        slide,
        Inches(7.0),
        Inches(3.9),
        Inches(5.7),
        Inches(2.75),
        "Atenção (QR / não oficial)",
        [
            "• Não é Cloud API da Meta.",
            "• Pode desconectar ou sofrer bloqueio.",
            "• Evite uso comercial massivo neste número.",
            "• Ideal só para bot da equipe interna.",
        ],
        title_color=WARNING,
        accent_bar=WARNING,
    )


def slide_conexao_oficial(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 2b — Conexão oficial (Número B / cliente)", "WhatsApp Oficial · Templates WABA")
    _card(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(6.2),
        Inches(5.5),
        "O que configurar (com a WhatsAtende)",
        [
            "1. WABA no Business Manager da ClickUp",
            "2. Número verificado na Meta",
            "3. No painel: Nova conexão → WhatsApp Oficial",
            "4. Preencher IDs (phone_number_id, waba_id, etc.)",
            "5. Copiar webhook Meta (waba_webhook) se pedido",
            "6. Criar templates (utilidade / marketing)",
            "7. Aguardar aprovação Meta",
            "",
            "Menu: Templates WABA",
            "Exemplos: lembrete de agenda, fatura, boas-vindas.",
        ],
        title_color=ACCENT_HOVER,
        accent_bar=ACCENT,
    )
    _card(
        slide,
        Inches(6.9),
        Inches(1.15),
        Inches(5.8),
        Inches(5.5),
        "Regra das 24 horas (simples)",
        [
            "Cliente falou com você → janela de 24h aberta.",
            "Dentro das 24h: texto/mídia livre via API.",
            "Fora das 24h: só template aprovado.",
            "",
            "Isso é regra da Meta, não da WhatsAtende.",
            "",
            "Para o bot interno (Número A / QR)",
            "essa trava não funciona da mesma forma.",
            "",
            "Por isso separamos os dois números.",
        ],
        title_color=TEAL,
    )


def slide_token(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 3 — Token da API (obrigatório)", "Sem token, o site-clickup não envia nada")
    _screen_mock(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(6.3),
        Inches(5.5),
        "Conexões · Editar conexão",
        [
            "1. Abra Conexões",
            "2. Clique em Editar na conexão do bot",
            "3. Localize o campo Token",
            "4. Gere / salve um token forte",
            "5. Copie o valor (key_conexao)",
            "",
            "Esse token vai no header:",
            "Authorization: Bearer SEU_TOKEN",
            "",
            "Guarde também o ID da conexão",
            "(whatsappId) — aparece nesta tela.",
        ],
        tip="Não compartilhe o token em print/chat público.",
    )
    _card(
        slide,
        Inches(7.0),
        Inches(1.15),
        Inches(5.7),
        Inches(5.5),
        "Para onde vai no nosso projeto",
        [
            "No Railway (produção) ou .env local:",
            "",
            "WHATSATENDE_API_URL=",
            "  https://api.app14.whatsatende.com.br",
            "",
            "WHATSATENDE_TOKEN=",
            "  <cole o token da conexão>",
            "",
            "WHATSATENDE_WHATSAPP_ID=",
            "  <cole o ID da conexão>",
            "",
            "Depois: redeploy / reinicie o serviço.",
        ],
        title_color=TEAL,
    )


def slide_api_teste(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 4 — Testar a API no painel", "Menu API → Teste API / Documentação API / Logs")
    _card(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Documentação API",
        [
            "Mostra todos os endpoints:",
            "• createConnection",
            "• startSession / getQrCode",
            "• checkStatus / disconnect",
            "• send (texto)",
            "• linkImage / linkPDF",
            "• base64",
            "• checkNumber",
            "• contacts / tickets",
            "",
            "Use como referência oficial.",
        ],
        title_color=TEAL,
    )
    _card(
        slide,
        Inches(4.7),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Teste API (primeiro envio)",
        [
            "1. Cole o Token",
            "2. Número: 55 + DDD + celular",
            "   Ex.: 5531999999999",
            "3. Body: Olá, teste ClickUp",
            "4. Enviar",
            "",
            "Se chegou no WhatsApp = token OK.",
            "",
            "Teste também PDF/imagem",
            "se for usar fatura/anexo.",
        ],
        title_color=ACCENT_HOVER,
        accent_bar=ACCENT,
    )
    _card(
        slide,
        Inches(9.0),
        Inches(1.15),
        Inches(3.8),
        Inches(5.5),
        "Logs API",
        [
            "Se der erro:",
            "• abra Logs API",
            "• confira status HTTP",
            "• token inválido?",
            "• número sem DDI?",
            "• conexão offline?",
            "",
            "Anote horário e payload",
            "para o suporte.",
        ],
        title_color=WARNING,
        accent_bar=WARNING,
    )


def slide_webhook(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 5 — Webhook (o bot “ouvir” mensagens)", "Crítico para o site-clickup funcionar como na Z-API")
    _card(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(6.2),
        Inches(5.5),
        "O que é e por que importa",
        [
            "Envio = site-clickup → WhatsAtende → WhatsApp.",
            "Recebimento = WhatsApp → WhatsAtende → nosso servidor.",
            "",
            "Sem webhook de mensagem recebida,",
            "o bot não responde comandos (VENDER, DFV, etc.).",
            "",
            "URL do nosso projeto:",
            "https://site-clickup-production.up.railway.app/api/crm/webhook-whatsapp/",
            "",
            "Confirme com o suporte:",
            "• como cadastrar a URL",
            "• payload JSON de exemplo",
            "• se precisa token/HMAC",
            "• se já está no plano app14",
        ],
        title_color=DANGER,
        accent_bar=DANGER,
    )
    _card(
        slide,
        Inches(6.9),
        Inches(1.15),
        Inches(5.8),
        Inches(5.5),
        "Status de entrega (também importante)",
        [
            "Webhook de status avisa:",
            "enviada / entregue / lida / falha.",
            "",
            "No site-clickup usamos isso para",
            "saber se o WhatsApp aceitou o envio.",
            "",
            "Enquanto o suporte não confirmar,",
            "dá para testar só o envio (outbound).",
            "O bot completo depende do inbound.",
            "",
            "Código já tem normalizer preparado",
            "para ajustar quando o payload chegar.",
        ],
        title_color=TEAL,
    )


def slide_site_record(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 6 — Ativar no site-clickup", "Painel interno ClickUp")
    _screen_mock(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(6.4),
        Inches(5.5),
        "ClickUp · Configuração WhatsApp",
        [
            "1. Entre na área interna (Diretoria/Admin/BackOffice)",
            "2. Abra Configuração WhatsApp",
            "3. Confira as linhas:",
            "   WhatsAtende: OK (conexão completa)",
            "4. Selecione o card WhatsAtende",
            "5. Clique em Salvar provedor",
            "6. (Opcional) Gerar QR / ver status no painel",
            "",
            "Enquanto as variáveis não estiverem",
            "no Railway, o card mostra 'ausente'.",
            "",
            "Padrão atual de produção: Z-API.",
            "Só troque quando o token/ID estiverem OK.",
        ],
        tip="Trocar provedor sem webhook inbound = bot mudo.",
    )
    _card(
        slide,
        Inches(7.1),
        Inches(1.15),
        Inches(5.6),
        Inches(5.5),
        "Variáveis Railway (produção)",
        [
            "WHATSATENDE_API_URL",
            "WHATSATENDE_TOKEN",
            "WHATSATENDE_WHATSAPP_ID",
            "",
            "Opcional:",
            "WHATSAPP_PROVIDER=whatsatende",
            "(ou deixe Z-API e troque só no painel)",
            "",
            "Migration 0187 já sobe o choice",
            "whatsatende no banco.",
            "",
            "Código já em produção (main).",
        ],
        title_color=SUCCESS,
        accent_bar=SUCCESS,
    )


def slide_teste_final(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Passo 7 — Teste ponta a ponta", "Checklist antes de usar com a equipe")
    _card(
        slide,
        Inches(0.4),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "1) Só envio (outbound)",
        [
            "☐ Conexão Conectada",
            "☐ Token no Railway",
            "☐ Provedor WhatsAtende",
            "☐ Enviar texto de teste",
            "☐ Enviar PDF/imagem",
            "☐ checkNumber OK",
            "",
            "Se isso passar,",
            "a API de envio está ok.",
        ],
        title_color=TEAL,
    )
    _card(
        slide,
        Inches(4.7),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "2) Bot completo (inbound)",
        [
            "☐ Webhook configurado",
            "☐ Mande 'oi' / comando teste",
            "☐ site-clickup recebe POST",
            "☐ Bot responde",
            "",
            "Pendências comuns:",
            "• botões / lista",
            "• payload diferente",
            "• webhook não incluso",
            "",
            "Aguarde respostas do suporte.",
        ],
        title_color=ACCENT_HOVER,
        accent_bar=ACCENT,
    )
    _card(
        slide,
        Inches(9.0),
        Inches(1.15),
        Inches(3.8),
        Inches(5.5),
        "3) Não fazer ainda",
        [
            "✗ Migrar tudo da Z-API",
            "✗ Usar QR para cliente final",
            "✗ Disparo massivo no QR",
            "✗ Expor token em prints",
            "✗ Trocar provedor sem teste",
            "",
            "Z-API continua plano B",
            "até paridade confirmada.",
        ],
        title_color=DANGER,
        accent_bar=DANGER,
    )


def slide_menu_rapido(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Onde clicar no menu (referência rápida)", "Telas que você mais vai usar")
    items = [
        ("Atendimentos", "Ver conversas/tickets humanos"),
        ("Conexões", "QR, status, token, ID da conexão"),
        ("API → Documentação", "Lista de endpoints e exemplos"),
        ("API → Teste API", "Enviar mensagem de prova"),
        ("API → Logs API", "Debug de erros de integração"),
        ("Templates WABA", "Templates oficiais (Número B)"),
        ("Integrações", "n8n / Typebot (alternativa inbound)"),
        ("Fluxos / Automações", "Chatbot interno da plataforma"),
        ("Configurações", "Preferências da empresa"),
    ]
    y = Inches(1.15)
    for i, (title, detail) in enumerate(items):
        col = i % 3
        row = i // 3
        left = Inches(0.4) + Inches(col * 4.25)
        top = y + Inches(row * 1.8)
        _card(
            slide,
            left,
            top,
            Inches(4.05),
            Inches(1.55),
            title,
            [detail],
            title_color=TEAL,
        )


def slide_resumo(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Resumo em 60 segundos", "O mínimo para lembrar")
    _card(
        slide,
        Inches(0.4),
        Inches(1.2),
        Inches(12.5),
        Inches(5.4),
        "Fórmula ClickUp + WhatsAtende",
        [
            "1. Login em app14 → Conexões → criar número do bot (QR).",
            "2. Gerar Token + copiar whatsappId.",
            "3. Colocar WHATSATENDE_TOKEN e WHATSATENDE_WHATSAPP_ID no Railway.",
            "4. Testar envio em API → Teste API.",
            "5. Configurar webhook inbound para /api/crm/webhook-whatsapp/ (com suporte).",
            "6. No site-clickup: Configuração WhatsApp → salvar provedor WhatsAtende.",
            "7. Testar comando do bot. Se não responder, o problema quase sempre é o webhook.",
            "8. Número do cliente final = Cloud API oficial (separado do bot interno).",
            "",
            "Documentos úteis na pasta docs/:",
            "• WhatsAtende_API_Documentacao.docx",
            "• Este PPT: WhatsAtende_Guia_Iniciante_Record_PAP.pptx",
        ],
        title_color=TEAL,
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)
    slide_objetivo(prs)
    slide_mapa(prs)
    slide_login(prs)
    slide_conexao_qr(prs)
    slide_conexao_oficial(prs)
    slide_token(prs)
    slide_api_teste(prs)
    slide_webhook(prs)
    slide_site_record(prs)
    slide_teste_final(prs)
    slide_menu_rapido(prs)
    slide_resumo(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
